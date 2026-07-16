"""제출용 PPT(4페이지) 생성.

결과가 바뀌면 다시 만들 수 있도록 문서를 코드로 남긴다.

실행:
    python make_report.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

WIDTH, HEIGHT = Inches(13.333), Inches(7.5)
FONT = "맑은 고딕"

INK = RGBColor(0x1F, 0x23, 0x28)
MUTED = RGBColor(0x65, 0x6D, 0x76)
ACCENT = RGBColor(0x0B, 0x5F, 0xFF)
GOOD = RGBColor(0x1A, 0x7F, 0x37)
BAD = RGBColor(0xCF, 0x22, 0x2E)
RULE = RGBColor(0xD8, 0xDE, 0xE4)
BAND = RGBColor(0xF4, 0xF6, 0xF8)

HERE = Path(__file__).parent
SAMPLES = HERE / "preprocessed_samples"
REPO_URL = "https://github.com/park-hyun-su/comento"


def add_slide(prs: Presentation, title: str, kicker: str) -> object:
    """제목 + 머리말 + 구분선만 놓인 빈 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(12), Inches(0.3))
    para = box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = kicker
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = ACCENT
    run.font.name = FONT

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.5))
    para = box.text_frame.paragraphs[0]
    run = para.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = INK
    run.font.name = FONT

    line = slide.shapes.add_connector(1, Inches(0.6), Inches(1.24), Inches(12.73), Inches(1.24))
    line.line.color.rgb = RULE
    line.line.width = Pt(1)
    return slide


def add_text(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, int, bool, RGBColor]],
    space_after: int = 6,
) -> None:
    """(텍스트, 크기, 굵게, 색) 목록을 한 텍스트박스에 쌓는다."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = FONT


def add_table(slide, left, top, width, height, rows: list[list[str]], col_widths=None):
    """헤더 1행 + 본문인 표."""
    shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r, row in enumerate(rows):
        table.rows[r].height = Inches(0.32)
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = BAND if r == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size = Pt(11)
                run.font.bold = r == 0
                run.font.name = FONT
                run.font.color.rgb = INK
                if r > 0 and c == len(row) - 1 and value.startswith("PASS"):
                    run.font.color.rgb = GOOD
                    run.font.bold = True
    return table


def slide1(prs: Presentation) -> None:
    slide = add_slide(prs, "Git 활용 및 픽셀 단위 이미지 처리", "1차 업무 · 요약")

    add_text(
        slide,
        Inches(0.6),
        Inches(1.5),
        Inches(6.1),
        Inches(4.5),
        [
            ("목표", 14, True, ACCENT),
            ("Git 브랜치·PR 기반 코드 관리 흐름을 익히고, OpenCV로 픽셀 단위", 13, False, INK),
            ("이미지 처리를 실습한다.", 13, False, INK),
            ("", 8, False, INK),
            ("수행 범위", 14, True, ACCENT),
            ("· 데이터: Hugging Face ethz/food101 무작위 5장 (seed 42)", 13, False, INK),
            ("· 전처리 10단계를 5장 각각에 전부 적용", 13, False, INK),
            ("· 기본 문제: 224×224, Grayscale, Normalize, Blur, 증강 3종", 13, False, INK),
            ("· 요청내용 2: HSV 기반 특정 색상 픽셀 감지 및 필터링", 13, False, INK),
            ("· 심화 문제: 어두운 이미지 / 객체 작은 이미지 이상치 탐지", 13, False, INK),
            ("", 8, False, INK),
            ("결과", 14, True, ACCENT),
            ("· 5장 전원 통과. 결과 이미지 preprocessed_samples/ 5장", 13, False, INK),
            ("· 브랜치 4개 → PR 4건 → main 병합", 13, False, INK),
            ("· cv2 연산을 NumPy로 재구현해 대조, 오차 ≤ 1 (반올림 수준)", 13, False, INK),
        ],
    )

    add_text(
        slide,
        Inches(7.0),
        Inches(1.5),
        Inches(5.7),
        Inches(0.4),
        [("데이터 흐름", 14, True, ACCENT)],
    )

    steps = [
        ("① 적재", "ethz/food101 · streaming · 무작위 5장"),
        ("② 기본 전처리", "224×224 → Grayscale → Normalize → Blur"),
        ("③ 데이터 증강", "좌우 반전 · 회전 +15° · HSV 색상 변화"),
        ("④ 색상 감지", "HSV inRange → 마스크 → 필터링"),
        ("⑤ 이상치 판정", "평균 밝기 · 최대 연결요소 면적 → PASS/REMOVE"),
    ]
    top = Inches(2.0)
    for name, detail in steps:
        box = slide.shapes.add_textbox(Inches(7.0), top, Inches(5.7), Inches(0.62))
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = name
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = INK
        run.font.name = FONT
        para2 = tf.add_paragraph()
        run2 = para2.add_run()
        run2.text = detail
        run2.font.size = Pt(11)
        run2.font.color.rgb = MUTED
        run2.font.name = FONT
        top = Emu(int(top) + int(Inches(0.78)))

    add_text(
        slide,
        Inches(0.6),
        Inches(6.55),
        Inches(12.1),
        Inches(0.4),
        [(f"저장소: {REPO_URL}", 12, True, ACCENT)],
    )


def slide2(prs: Presentation) -> None:
    slide = add_slide(prs, "Git 워크플로우 — 브랜치 · PR · 병합", "1차 업무 · 코드 관리")

    add_text(
        slide,
        Inches(0.6),
        Inches(1.5),
        Inches(6.1),
        Inches(0.4),
        [("기능 단위 브랜치 전략", 14, True, ACCENT)],
    )

    branches = [
        ("main", "보호 대상. PR로만 병합", MUTED),
        ("feature/image-processing", "기본 전처리 + NumPy 대조 검증", INK),
        ("feature/color-detection", "HSV 색상 감지 / 필터링", INK),
        ("feature/outlier-filter", "이상치 탐지 (어두움 / 작은 객체)", INK),
        ("feature/pipeline-report", "파이프라인 · 결과 이미지 · 문서", INK),
    ]
    top = Inches(1.95)
    for name, detail, color in branches:
        box = slide.shapes.add_textbox(Inches(0.6), top, Inches(6.1), Inches(0.55))
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        run = para.add_run()
        run.text = name
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = color
        run.font.name = "Consolas"
        run2 = para.add_run()
        run2.text = f"   {detail}"
        run2.font.size = Pt(11)
        run2.font.color.rgb = MUTED
        run2.font.name = FONT
        top = Emu(int(top) + int(Inches(0.62)))

    add_text(
        slide,
        Inches(0.6),
        Inches(5.2),
        Inches(6.1),
        Inches(1.6),
        [
            ("커밋 규칙", 14, True, ACCENT),
            ("· 한 커밋에 한 가지 변경만 담는다", 12, False, INK),
            ("· 제목은 타입(feat/fix/docs) + 무엇을 바꿨는지", 12, False, INK),
            ("· 본문에는 '왜'와 검증 결과를 남긴다", 12, False, INK),
        ],
    )

    add_text(
        slide,
        Inches(7.0),
        Inches(1.5),
        Inches(5.7),
        Inches(0.4),
        [("사용한 명령", 14, True, ACCENT)],
    )

    box = slide.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.7), Inches(3.4))
    box.fill.solid()
    box.fill.fore_color.rgb = BAND
    box.line.color.rgb = RULE
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_top = Inches(0.14)
    commands = [
        "# 저장소 연동",
        "git clone https://github.com/park-hyun-su/comento.git",
        "",
        "# 브랜치 생성 및 이동",
        "git checkout -b feature/image-processing",
        "",
        "# 변경 사항 관리",
        "git add . && git commit -m \"feat: ...\"",
        "git push -u origin feature/image-processing",
        "",
        "# PR 생성 → 리뷰 → 병합",
        "gh pr create --base main",
        "git checkout main && git merge feature/image-processing",
        "git push origin main",
    ]
    for i, line in enumerate(commands):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(1)
        run = para.add_run()
        run.text = line
        run.font.size = Pt(10.5)
        run.font.name = "Consolas"
        run.font.color.rgb = MUTED if line.startswith("#") else INK

    add_text(
        slide,
        Inches(7.0),
        Inches(5.55),
        Inches(5.7),
        Inches(1.2),
        [
            ("리뷰 포인트로 올린 것", 14, True, ACCENT),
            ("PR 본문에 cv2 대조 검증 수치와, 구현 중 드러난", 12, False, INK),
            ("두 가지 함정(작은 커널 분기 · 회전 행렬 방향)을 기록", 12, False, INK),
        ],
    )


def slide3(prs: Presentation) -> None:
    slide = add_slide(prs, "전처리 파이프라인 — 10단계", "1차 업무 · 픽셀 단위 이미지 처리")

    # 몽타주는 20x9.2 비율이라 폭을 그대로 12.13"로 두면 높이가 5.6"까지 나와
    # 아래 텍스트를 덮는다. 폭 8.5" -> 높이 3.91"로 맞춰 좌측에 배치한다.
    sample = SAMPLES / "1_peking_duck_steps.png"
    if sample.exists():
        slide.shapes.add_picture(str(sample), Inches(0.6), Inches(1.5), width=Inches(8.5))

    add_text(
        slide,
        Inches(0.6),
        Inches(5.55),
        Inches(8.5),
        Inches(1.7),
        [
            ("단계", 13, True, ACCENT),
            ("① 224×224 (INTER_AREA)  ② Grayscale (BT.601)  ③ Normalize (ImageNet)  ④ Blur 5×5  ⑤ 좌우 반전", 11.5, False, INK),
            ("⑥ 회전 +15° (반사 패딩)  ⑦ HSV 색상 변화  ⑧ Red 마스크  ⑨ Red 픽셀만  ⑩ 객체 마스크", 11.5, False, INK),
            ("5장 모두 같은 형식으로 preprocessed_samples/에 저장", 10.5, False, MUTED),
        ],
    )

    add_text(
        slide,
        Inches(9.3),
        Inches(1.5),
        Inches(3.45),
        Inches(5.4),
        [
            ("왜 RGB가 아니라", 13, True, ACCENT),
            ("HSV에서 색을 고르는가", 13, True, ACCENT),
            ("", 3, False, INK),
            ("RGB는 밝기가 바뀌면 임계값이", 11.5, False, INK),
            ("같이 흔들린다. HSV는 색상(H)이", 11.5, False, INK),
            ("밝기(V)와 분리돼 조명이 달라져도", 11.5, False, INK),
            ("같은 H 범위로 잡힌다.", 11.5, False, INK),
            ("", 3, False, INK),
            ("다만 두 가지를 처리해야 했다", 11.5, True, INK),
            ("", 3, False, INK),
            ("· 빨강은 H=0에서 갈라진다.", 11.5, False, INK),
            ("  0~10과 170~179를 OR로 합침", 11.5, False, INK),
            ("", 3, False, INK),
            ("· 회색 접시는 HSV=(0,0,128).", 11.5, False, INK),
            ("  H가 빨강과 똑같은 0이다.", 11.5, False, BAD),
            ("  채도(S) 하한을 둬야 접시가", 11.5, False, INK),
            ("  빨강으로 안 잡힌다.", 11.5, False, INK),
            ("  (합성 이미지로 누출 0px 확인)", 11.5, False, GOOD),
        ],
        space_after=2,
    )


def slide4(prs: Presentation) -> None:
    slide = add_slide(prs, "이상치 탐지 및 검증", "1차 업무 · 심화 문제 + 결과")

    add_text(
        slide,
        Inches(0.6),
        Inches(1.42),
        Inches(6.1),
        Inches(0.4),
        [("무작위 5장 판정 결과 (seed 42)", 13, True, ACCENT)],
    )
    add_table(
        slide,
        Inches(0.6),
        Inches(1.8),
        Inches(6.1),
        Inches(2.0),
        [
            ["라벨", "평균 밝기", "객체 면적", "판정"],
            ["peking_duck", "110.0", "48.6%", "PASS"],
            ["tiramisu", "127.5", "30.3%", "PASS"],
            ["beignets", "98.6", "66.7%", "PASS"],
            ["ravioli", "139.1", "42.5%", "PASS"],
            ["tacos", "121.4", "51.8%", "PASS"],
        ],
        col_widths=[Inches(2.2), Inches(1.4), Inches(1.4), Inches(1.1)],
    )

    add_text(
        slide,
        Inches(0.6),
        Inches(4.02),
        Inches(6.1),
        Inches(0.4),
        [("5장이 전부 정상이라, 제거 분기는 합성 이미지로 검증", 13, True, ACCENT)],
    )
    add_text(
        slide,
        Inches(0.6),
        Inches(6.25),
        Inches(6.1),
        Inches(0.8),
        [
            ("값이 균일한 이미지는 Otsu가 나눌 경계가 없어 면적이 100%로 나온다.", 10.5, False, MUTED),
            ("밝기 검사가 먼저 걸러내므로 판정은 옳지만 이 수치 자체엔 의미가 없다.", 10.5, False, MUTED),
        ],
        space_after=2,
    )
    add_table(
        slide,
        Inches(0.6),
        Inches(4.4),
        Inches(6.1),
        Inches(1.7),
        [
            ["케이스", "면적", "이론값", "판정"],
            ["균일한 어두운 이미지", "100%", "퇴화 사례", "어두움 제거"],
            ["밝은 배경 · 20×20 객체", "0.8%", "0.8%", "객체 작음 제거"],
            ["밝은 배경 · 140×140 객체", "39.0%", "39.1%", "PASS"],
            ["어두운 배경 · 밝은 객체", "39.0%", "39.1%", "PASS"],
        ],
        col_widths=[Inches(2.6), Inches(1.1), Inches(1.1), Inches(1.3)],
    )

    add_text(
        slide,
        Inches(7.0),
        Inches(1.42),
        Inches(5.73),
        Inches(2.9),
        [
            ("전경을 고르는 문제 — 실제로 틀렸던 것", 13, True, ACCENT),
            ("Otsu는 밝기로 두 덩어리를 나눌 뿐 어느 쪽이 피사체인지 모른다.", 11.5, False, INK),
            ("접시가 밝은 사진과 어두운 사진이 섞여 극성을 고정할 수 없다.", 11.5, False, INK),
            ("", 3, False, INK),
            ("처음엔 '테두리를 덜 차지하는 쪽'을 전경으로 골랐다. beignets는", 11.5, False, INK),
            ("음식이 화면을 꽉 채워 테두리를 많이 건드린 탓에 배경으로 오인돼,", 11.5, False, INK),
            ("배경 조각을 잡고 15.7%를 보고했다. (PASS는 우연이었다)", 11.5, False, BAD),
            ("", 3, False, INK),
            ("전경다움 = 중심 커버리지 − 테두리 점유 → 66.7%로 정정", 11.5, True, GOOD),
            ("한쪽만 쓰면 깨진다. 테두리만 → 꽉 찬 피사체가 배경 취급.", 11.5, False, MUTED),
            ("중심만 → 작은 피사체를 감싼 배경이 이긴다.", 11.5, False, MUTED),
        ],
        space_after=3,
    )

    add_text(
        slide,
        Inches(7.0),
        Inches(4.4),
        Inches(5.73),
        Inches(2.7),
        [
            ("NumPy 직접 구현 대조 — 오차 ≤ 1 (반올림 수준)", 13, True, ACCENT),
            ("grayscale 0.0011 · blur 0.0021 · flip 0.0000 · rotate 0.0004", 11.5, False, INK),
            ("(cv2 5.0 기준. 4.12는 회전만 0.81 — warpAffine이 5비트 고정소수점", 10.5, False, MUTED),
            ("보간을 쓴다. 전처리 결과와 판정은 두 버전 동일)", 10.5, False, MUTED),
            ("", 3, False, INK),
            ("대조하지 않았으면 못 찾았을 두 가지", 11.5, True, INK),
            ("· cv2는 sigma≤0 & ksize≤7이면 공식이 아니라 하드코딩 이항계수", 11.5, False, INK),
            ("  테이블을 쓴다. 공식대로 짜면 오차 0.73이 남는다.", 11.5, False, INK),
            ("· getRotationMatrix2D는 순방향 행렬. 역방향 매핑에 그대로 쓰면", 11.5, False, INK),
            ("  반대로 회전한다 (오차 57.8 → 역행렬로 0.0004).", 11.5, False, INK),
            ("", 3, False, INK),
            ("한계: Otsu 기반이라 진짜 객체 탐지가 아닌 면적 프록시다.", 11.5, False, MUTED),
            ("tacos는 밝은 접시·밥은 잡고 어두운 콩은 뺀다.", 11.5, False, MUTED),
        ],
        space_after=3,
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = WIDTH, HEIGHT

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)

    out = HERE / "1차업무_결과보고_박현수.pptx"
    prs.save(out)
    print(f"저장 완료: {out} ({len(prs.slides._sldIdLst)}장)")


if __name__ == "__main__":
    main()
