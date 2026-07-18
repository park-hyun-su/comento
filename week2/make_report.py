"""2차 업무 제출용 PPT(4페이지) 생성 — 1주차 make_report.py와 동일한 스타일.

실행:
    python make_report.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

WIDTH, HEIGHT = Inches(13.333), Inches(7.5)
FONT = "맑은 고딕"

INK = RGBColor(0x1F, 0x23, 0x28)
MUTED = RGBColor(0x65, 0x6D, 0x76)
ACCENT = RGBColor(0x0B, 0x5F, 0xFF)
GOOD = RGBColor(0x1A, 0x7F, 0x37)
BAD = RGBColor(0xCF, 0x22, 0x2E)
RULE = RGBColor(0xD8, 0xDE, 0xE4)
BAND = RGBColor(0xF4, 0xF6, 0xF8)

HERE = Path(__file__).parent
OUT = HERE / "outputs"
REPO_URL = "https://github.com/park-hyun-su/comento"


def add_slide(prs, title, kicker):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(12), Inches(0.3))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = kicker
    run.font.size, run.font.bold, run.font.color.rgb, run.font.name = Pt(11), True, ACCENT, FONT

    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(12), Inches(0.5))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = title
    run.font.size, run.font.bold, run.font.color.rgb, run.font.name = Pt(26), True, INK, FONT

    line = slide.shapes.add_connector(1, Inches(0.6), Inches(1.24), Inches(12.73), Inches(1.24))
    line.line.color.rgb, line.line.width = RULE, Pt(1)
    return slide


def add_text(slide, left, top, width, height, lines, space_after=6):
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = text
        run.font.size, run.font.bold, run.font.color.rgb, run.font.name = Pt(size), bold, color, FONT


def add_table(slide, left, top, width, height, rows, col_widths=None):
    table = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows):
        table.rows[r].height = Inches(0.3)
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = BAND if r == 0 else RGBColor(0xFF, 0xFF, 0xFF)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            for run in para.runs:
                run.font.size, run.font.name, run.font.color.rgb = Pt(10.5), FONT, INK
                run.font.bold = r == 0


def slide1(prs):
    slide = add_slide(prs, "Unit Test 구성 및 2D → 3D 변환", "2차 업무 · 요약")
    add_text(slide, Inches(0.6), Inches(1.5), Inches(6.1), Inches(4.8), [
        ("목표", 14, True, ACCENT),
        ("pytest로 픽셀 처리·3D 변환 코드를 검증하고, OpenCV+NumPy로", 13, False, INK),
        ("2D 이미지를 Depth Map → 3D 포인트클라우드로 변환한다.", 13, False, INK),
        ("", 8, False, INK),
        ("수행 범위", 14, True, ACCENT),
        ("· Unit Test 16개 (그레이스케일·정규화·Depth·포인트클라우드·PLY)", 13, False, INK),
        ("· 2D→3D: 밝기를 높이(Z)로 보는 height-field 변환", 13, False, INK),
        ("· 합성 장면(정답 형상) + 실제 이미지 두 경로로 검증", 13, False, INK),
        ("· 결과 시각화: 원본·Depth·3D표면·포인트클라우드 4분할 + .ply", 13, False, INK),
        ("", 8, False, INK),
        ("결과", 14, True, ACCENT),
        ("· pytest 16개 전원 통과 (0.88s)", 13, False, GOOD),
        ("· 합성 반구+램프 형상 정확히 복원 → 알고리즘 검증", 13, False, INK),
        ("· 실제 이미지에서 '밝기≠깊이' 한계 확인 → 개선점 도출", 13, False, INK),
    ])
    add_text(slide, Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.4),
             [("파이프라인", 14, True, ACCENT)])
    steps = [
        ("① 입력 2D 이미지", "BGR / Grayscale"),
        ("② 그레이스케일 + 평활화", "cvtColor → GaussianBlur 5×5"),
        ("③ Depth Map", "normalize [0,1] · invert 옵션"),
        ("④ 3D 포인트클라우드", "X=col, Y=row, Z=depth×scale"),
        ("⑤ 저장 · 시각화", "PLY + matplotlib 3D"),
    ]
    top = Inches(2.0)
    for name, detail in steps:
        tf = slide.shapes.add_textbox(Inches(7.0), top, Inches(5.7), Inches(0.62)).text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = name
        r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(13), True, INK, FONT
        r2 = tf.add_paragraph().add_run()
        r2.text = detail
        r2.font.size, r2.font.color.rgb, r2.font.name = Pt(11), MUTED, "Consolas"
        top = Emu(int(top) + int(Inches(0.82)))
    add_text(slide, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
             [(f"저장소: {REPO_URL}", 12, True, ACCENT)])


def slide2(prs):
    slide = add_slide(prs, "Unit Test — 수치 정답으로 로직 검증", "2차 업무 · 코드 검증")
    add_text(slide, Inches(0.6), Inches(1.45), Inches(6.2), Inches(0.4),
             [("검증한 함수와 케이스", 13, True, ACCENT)])
    add_table(slide, Inches(0.6), Inches(1.85), Inches(6.2), Inches(3.0), [
        ["함수", "핵심 검증"],
        ["to_grayscale", "BT.601 known=141 · None→예외"],
        ["normalize01", "[0,1] · 균일영상 0나눗셈 방지"],
        ["estimate_depth", "범위 · invert 극성 뒤집힘"],
        ["depth_to_pointcloud", "Z=depth×scale 좌표대응 · 개수"],
        ["pointcloud_colors", "점–색 1:1 · BGR→RGB"],
        ["save_ply", "PLY 헤더·정점 개수"],
        ["backproject_pinhole", "주점 X=Y=0, Z=depth"],
    ], col_widths=[Inches(2.5), Inches(3.7)])
    add_text(slide, Inches(0.6), Inches(5.5), Inches(6.2), Inches(1.6), [
        ("멘토 예제와의 차이", 13, True, ACCENT),
        ("예제 테스트는 applyColorMap 출력의 shape/type만 본다 —", 11.5, False, INK),
        ("'죽지 않는다'만 보장한다. 여기서는 수치 정답이 있는 입력으로", 11.5, False, INK),
        ("변환 공식 자체를 못 박았다.", 11.5, False, INK),
    ], space_after=3)

    add_text(slide, Inches(7.0), Inches(1.45), Inches(5.7), Inches(0.4),
             [("pytest 실행 결과", 13, True, ACCENT)])
    box = slide.shapes.add_textbox(Inches(7.0), Inches(1.85), Inches(5.73), Inches(3.6))
    box.fill.solid()
    box.fill.fore_color.rgb = BAND
    box.line.color.rgb = RULE
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left, tf.margin_top = Inches(0.16), Inches(0.12)
    lines = [
        "$ python -m pytest -v",
        "collected 16 items",
        "",
        "test_to_grayscale_shape_and_dtype ....... PASSED",
        "test_to_grayscale_known_value ........... PASSED",
        "test_to_grayscale_none_raises ........... PASSED",
        "test_normalize01_uniform_is_zero ........ PASSED",
        "test_estimate_depth_invert_flips ........ PASSED",
        "test_pointcloud_z_equals_depth_scale .... PASSED",
        "test_pointcloud_count_subsampled ........ PASSED",
        "test_save_ply_header_and_count .......... PASSED",
        "test_backproject_pinhole_principal ...... PASSED",
        "  ... (16 tests)",
        "",
        "===== 16 passed in 0.88s =====",
    ]
    for i, ln in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(1)
        run = para.add_run()
        run.text = ln
        run.font.size, run.font.name = Pt(10), "Consolas"
        run.font.color.rgb = GOOD if ("PASSED" in ln or "passed" in ln) else (
            MUTED if ln.startswith(("$", "collected", "  ...")) else INK)
    add_text(slide, Inches(7.0), Inches(5.6), Inches(5.7), Inches(1.4), [
        ("대표 케이스", 13, True, ACCENT),
        ("test_pointcloud_z_equals_depth_times_scale:", 11, True, INK),
        ("(열,행)별 Z가 depth×z_scale와 정확히 일치하는지", 11, False, MUTED),
        ("좌표 단위로 확인 → 변환 공식이 맞음을 보증", 11, False, MUTED),
    ], space_after=2)


def slide3(prs):
    slide = add_slide(prs, "2D → 3D 변환 결과", "2차 업무 · 결과")
    for path, x, cap in [
        (OUT / "synthetic_panels.png", Inches(0.55), "합성 장면 (반구+램프) — 정답 형상 복원 검증"),
        (OUT / "tiramisu_photo_panels.png", Inches(6.95), "실제 이미지 (티라미수) — 생크림↑ 초콜릿↓"),
    ]:
        if path.exists():
            slide.shapes.add_picture(str(path), x, Inches(1.5), width=Inches(5.85))
        add_text(slide, x, Inches(6.35), Inches(5.85), Inches(0.6),
                 [(cap, 11.5, True, INK)])
    add_text(slide, Inches(0.55), Inches(6.9), Inches(12), Inches(0.4), [
        ("각 4분할: 1) 입력 2D  2) Depth map(JET)  3) 3D 표면(Z=depth)  4) 컬러 포인트클라우드  ·  결과 .ply는 MeshLab/CloudCompare/Open3D로 열람",
         10, False, MUTED)])


def slide4(prs):
    slide = add_slide(prs, "한계 · 개선점 및 참고 논문", "2차 업무 · 분석")
    add_text(slide, Inches(0.6), Inches(1.42), Inches(6.2), Inches(0.4),
             [("한계 → 개선 방향 (밝기를 깊이로 가정)", 13, True, ACCENT)])
    add_table(slide, Inches(0.6), Inches(1.82), Inches(6.2), Inches(2.2), [
        ["한계", "개선 방향"],
        ["밝기≠깊이 (초콜릿이 구멍으로)", "스테레오 / 학습기반 depth"],
        ["조명 의존 (그림자=가짜 높이)", "알베도-조명 분리"],
        ["절대 스케일 없음", "스테레오 캘리브 / metric depth"],
        ["노이즈 민감", "bilateral 등 엣지보존 평활화"],
    ], col_widths=[Inches(3.2), Inches(3.0)])
    add_text(slide, Inches(0.6), Inches(4.35), Inches(6.2), Inches(2.6), [
        ("다음 단계 두 경로", 13, True, ACCENT),
        ("① 스테레오: StereoSGBM → disparity →", 11.5, False, INK),
        ("   reprojectImageTo3D(Q).  backproject_pinhole()이", 11.5, False, INK),
        ("   이 경로의 역투영을 미리 구현.", 11.5, False, MUTED),
        ("② 단일 이미지 딥러닝 depth: MiDaS / Depth Anything V2를", 11.5, False, INK),
        ("   붙이면 캘리브 없이 진짜 상대깊이 → 같은", 11.5, False, INK),
        ("   depth_to_pointcloud()로 포인트클라우드 생성.", 11.5, False, MUTED),
    ], space_after=3)

    add_text(slide, Inches(7.0), Inches(1.42), Inches(5.73), Inches(5.5), [
        ("참고 논문 — 2D → 3D 변환", 13, True, ACCENT),
        ("단일 이미지 → 깊이 (본 과제와 최근접)", 11.5, True, INK),
        ("· Eigen+ 2014 — 단일영상 depth 회귀의 시초", 10.5, False, INK),
        ("· MiDaS (Ranftl+ 2022) — robust 상대깊이", 10.5, False, INK),
        ("· Depth Anything V2 (2024) — depth 파운데이션", 10.5, False, INK),
        ("· Horn 1970 — Shape from Shading (밝기→높이)", 10.5, False, GOOD),
        ("", 4, False, INK),
        ("스테레오 / disparity (OpenCV 근거)", 11.5, True, INK),
        ("· Scharstein & Szeliski 2002 — 스테레오 taxonomy", 10.5, False, INK),
        ("· Hirschmüller 2008 — SGM (StereoSGBM 내부)", 10.5, False, INK),
        ("", 4, False, INK),
        ("다중뷰 · 신경망 3D", 11.5, True, INK),
        ("· COLMAP 2016 · NeRF 2020 · 3DGS 2023", 10.5, False, INK),
        ("· Zero-1-to-3 2023 · TripoSR 2024 (단일→3D 생성)", 10.5, False, INK),
        ("", 4, False, INK),
        ("위성 2D→3D (DSM/DEM)", 11.5, True, ACCENT),
        ("· IM2HEIGHT 2018 · IM2ELEVATION 2020", 10.5, False, INK),
        ("· S2P 2014 (위성 스테레오) · IEEE GRSS DFC'23", 10.5, False, INK),
    ], space_after=2)


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = WIDTH, HEIGHT
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    out = HERE / "2차업무_결과보고_박현수.pptx"
    prs.save(out)
    print(f"저장 완료: {out} ({len(prs.slides._sldIdLst)}장)")


if __name__ == "__main__":
    main()
